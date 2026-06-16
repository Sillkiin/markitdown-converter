import os, sys, json, zipfile, tempfile, subprocess, importlib.util, types, struct, wave, math
SKILL=os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
CONVERT=os.path.join(SKILL,"scripts","convert.py")
spec=importlib.util.spec_from_file_location("convert",CONVERT); conv=importlib.util.module_from_spec(spec); spec.loader.exec_module(conv)
ns=types.SimpleNamespace(use_plugins=False,docintel_endpoint=None,llm_model=None,llm_api_key=None,llm_base_url=None)
md=conv.build_converter(ns)
results=[]; stderr_blobs=[]
def check(n,c,d=""): results.append((n,bool(c),d))
def ct(p): return conv.convert_one(md,p)
def cli(a,cwd=None):
    pr=subprocess.run([sys.executable,CONVERT,*a],capture_output=True,text=True,cwd=cwd); stderr_blobs.append(pr.stderr); return pr.returncode,pr.stdout,pr.stderr
D=tempfile.mkdtemp()
def P(n): return os.path.join(D,n)

# ----- build samples -----
json.dump({"product":"Виджет","price":42,"nested":{"a":[1,2,3]}},open(P("data.json"),"w",encoding="utf-8"),ensure_ascii=False)
open(P("table.csv"),"w").write("name,role\nNik,founder\nIvan,dev\n")
open(P("uni.csv"),"w",encoding="utf-8").write("город,страна\nМосква,Россия\n")
open(P("data.tsv"),"w").write("a\tb\n1\t2\n")
open(P("empty.csv"),"w").write("")
open(P("page.html"),"w").write("<h1>Title</h1><p>Some <b>bold</b> and <a href='https://x.com'>link</a>.</p><ul><li>a</li><li>b</li></ul>")
open(P("entities.html"),"w").write("<p>caf&eacute; &amp; co &mdash; 5&lt;6</p>")
open(P("feed.xml"),"w").write('<?xml version="1.0"?><catalog><item><name>Pen</name></item></catalog>')
open(P("feed.rss"),"w").write('<?xml version="1.0"?><rss version="2.0"><channel><title>News</title><item><title>Item A</title><description>desc A</description></item></channel></rss>')
json.dump({"cells":[{"cell_type":"markdown","metadata":{},"source":["# Title\n","text"]},{"cell_type":"code","metadata":{},"execution_count":1,"outputs":[],"source":["print('hi')"]}],"metadata":{},"nbformat":4,"nbformat_minor":5}, open(P("nb.ipynb"),"w"))
with zipfile.ZipFile(P("bundle.zip"),"w") as z: z.writestr("inner.csv","a,b\n1,2\n"); z.writestr("note.html","<h2>Note</h2>")
with zipfile.ZipFile(P("nested.zip"),"w") as z: z.write(P("bundle.zip"),"bundle.zip"); z.writestr("top.csv","x\n9\n")
# epub
with zipfile.ZipFile(P("book.epub"),"w") as z:
    z.writestr("mimetype","application/epub+zip")
    z.writestr("META-INF/container.xml",'<?xml version="1.0"?><container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>')
    z.writestr("OEBPS/content.opf",'<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="2.0" unique-identifier="id"><metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>My Book</dc:title><dc:creator>Nik</dc:creator></metadata><manifest><item id="c1" href="ch1.xhtml" media-type="application/xhtml+xml"/></manifest><spine><itemref idref="c1"/></spine></package>')
    z.writestr("OEBPS/ch1.xhtml",'<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><body><h1>Chapter 1</h1><p>Hello epub body.</p></body></html>')
import openpyxl
wb=openpyxl.Workbook(); ws=wb.active; ws.title="Q2"; ws.append(["Metric","Value"]); ws.append(["MRR","$420k"]); w3=wb.create_sheet("Q3"); w3.append(["Metric","Value"]); w3.append(["MRR","$510k"]); wb.save(P("book.xlsx"))
from pptx import Presentation
pr=Presentation(); s=pr.slides.add_slide(pr.slide_layouts[1]); s.shapes.title.text="Roadmap"; s.placeholders[1].text_frame.text="Q1 beta"; pr.save(P("deck.pptx"))
from docx import Document
doc=Document(); doc.add_heading("Quarterly Report",1); doc.add_paragraph("Revenue grew.")
for it in ["alpha","beta"]: doc.add_paragraph(it,style="List Bullet")
t=doc.add_table(rows=2,cols=2); t.cell(0,0).text="Metric"; t.cell(0,1).text="Value"; t.cell(1,0).text="MRR"; t.cell(1,1).text="$1M"; doc.save(P("report.docx"))
from reportlab.pdfgen import canvas
c=canvas.Canvas(P("file.pdf")); c.drawString(72,720,"Hello PDF Marker"); c.drawString(72,700,"Second line"); c.save()
from PIL import Image, PngImagePlugin
im=Image.new("RGB",(40,20),(123,222,64)); meta=PngImagePlugin.PngInfo(); meta.add_text("Title","Sample PNG"); meta.add_text("Author","Nik"); im.save(P("pic.png"),pnginfo=meta)
Image.new("RGB",(30,30),(10,20,30)).save(P("pic.jpg"),"JPEG")
w=wave.open(P("a.wav"),"w"); w.setnchannels(1); w.setsampwidth(2); w.setframerate(8000); w.writeframes(b"".join(struct.pack("<h",int(1000*math.sin(i/10))) for i in range(8000))); w.close()
open(P("song.mp3"),"wb").write(b"ID3\x03\x00\x00\x00\x00\x00\x21"+b"\x00"*200)  # minimal mp3-ish
open(P("отчёт за Q2.csv"),"w",encoding="utf-8").write("k,v\nузел,значение\n")

# ----- content tests (in-process) -----
check("01 pdf","Hello PDF Marker" in ct(P("file.pdf")))
m=ct(P("report.docx")); check("02 docx heading","# Quarterly Report" in m); check("03 docx table","| MRR | $1M |" in m); check("04 docx bullets",("- alpha" in m or "* alpha" in m) and ("- beta" in m or "* beta" in m))
m=ct(P("book.xlsx")); check("05 xlsx table","$420k" in m); check("06 xlsx multi-sheet","$510k" in m)
m=ct(P("deck.pptx")); check("07 pptx","Roadmap" in m and "Q1 beta" in m)
check("08 csv table","| Nik | founder |" in ct(P("table.csv")))
check("09 tsv","| a | b |" in ct(P("data.tsv")) or "a\tb" in ct(P("data.tsv")))
m=ct(P("page.html")); check("10 html bold+link","**bold**" in m and "[link](https://x.com)" in m); check("11 html heading","# Title" in m)
check("12 html entities","café" in ct(P("entities.html")))
check("13 json unicode","Виджет" in ct(P("data.json")))
check("14 xml","Pen" in ct(P("feed.xml")))
check("15 rss","Item A" in ct(P("feed.rss")))
check("16 ipynb","# Title" in ct(P("nb.ipynb")) and "print('hi')" in ct(P("nb.ipynb")))
m=ct(P("bundle.zip")); check("17 zip iterate","inner.csv" in m and "Note" in m)
check("18 nested zip","top.csv" in ct(P("nested.zip")))
check("19 epub","Chapter 1" in ct(P("book.epub")) and "My Book" in ct(P("book.epub")))
check("20 csv cyrillic","Москва" in ct(P("uni.csv")))
m=ct(P("pic.png")); check("21 png metadata fallback","Image:" in m and "40 x 20" in m and "Sample PNG" in m)
m=ct(P("pic.jpg")); check("22 jpg metadata fallback","Image:" in m and "30 x 30" in m)
m=ct(P("a.wav")); check("23 wav metadata fallback","Audio:" in m and "8000 Hz" in m and "Duration:" in m)
m=ct(P("song.mp3")); check("24 mp3 note fallback (non-empty)", bool(m.strip()) and "MP3" in m.upper())
check("25 empty csv no crash", isinstance(ct(P("empty.csv")),str))
check("26 spaced+cyrillic filename","значение" in ct(P("отчёт за Q2.csv")))

# ----- CLI-mode tests -----
cdir=P("coll"); os.makedirs(cdir,exist_ok=True)
for e,c in [("csv","p,q\n1,2\n"),("html","<h3>H</h3>"),("json",'{"k":1}')]: open(os.path.join(cdir,f"same.{e}"),"w").write(c)
o=P("od"); rc,_,_=cli([cdir+"/","-o",o+"/"]); fs=sorted(os.listdir(o)); check("27 dir no collision",rc==0 and fs==["same.csv.md","same.html.md","same.json.md"],str(fs))
o=P("om"); rc,_,_=cli([P("data.json"),P("feed.xml"),P("file.pdf"),"-o",o+"/"]); check("28 multi args=3",rc==0 and len(os.listdir(o))==3)
o=P("os1"); rc,_,_=cli([P("book.xlsx"),"-o",o+"/"]); check("29 single->dir basename",rc==0 and os.path.exists(os.path.join(o,"book.md")))
outd=P("std"); os.makedirs(outd,exist_ok=True); rc,so,_=cli([P("table.csv")],cwd=outd); check("30 stdout+cwd file",rc==0 and "| name | role |" in so and os.path.exists(os.path.join(outd,"table.md")))
rc,_,er=cli([P("ghost.pdf"),"-o",P("g.md")]); check("31 nonexistent->exit2 clean",rc==2 and "FAILED" in er)
o=P("op"); rc,_,er=cli([P("data.json"),P("ghost.docx"),P("feed.xml"),"-o",o+"/"]); check("32 partial 2 ok exit2",rc==2 and len(os.listdir(o))==2)
rc,_,_=cli([P("page.html"),"--use-plugins","-o",P("pl.md")]); check("33 --use-plugins ok",rc==0 and "**bold**" in open(P("pl.md")).read())
# mixed-format directory end-to-end
mdir=P("mixed"); os.makedirs(mdir,exist_ok=True)
for f in ["file.pdf","report.docx","book.xlsx","deck.pptx","pic.png","a.wav","nb.ipynb","book.epub","feed.rss"]:
    import shutil; shutil.copy(P(f), os.path.join(mdir, f))
o=P("mixout"); rc,_,er=cli([mdir+"/","-o",o+"/"]); check("34 mixed dir all convert",rc==0 and len(os.listdir(o))==9, "%d files rc=%d"%(len(os.listdir(o)),rc))

import openpyxl as _ox
_wb=_ox.Workbook(); _ws=_wb.active; _ws.append(["Region","Q1"]); _ws.append(["NA",120]); _ws.append(["N/A",5]); _ws.append(["EU",80]); _wb.save(P("na.xlsx"))
_m=ct(P("na.xlsx")); check("36 xlsx literal NA preserved (no pandas NaN)", "| NA |" in _m and "| N/A |" in _m and "NaN" not in _m)

_wb2=_ox.Workbook(); _ws2=_wb2.active; _ws2.append(["Note","Detail"]); _ws2.append(["a|b","line1\nline2"]); _wb2.save(P("esc.xlsx"))
_m2=ct(P("esc.xlsx")); check("37 xlsx escapes pipe+newline (table not broken)", "a\\|b" in _m2 and "line1<br>line2" in _m2 and "NaN" not in _m2)
_wb3=_ox.Workbook(); _ws3=_wb3.active; _ws3.append(["x","y","total"]); _ws3.append([3,10,None]); _ws3["C2"]="=A2+B2"; _wb3.save(P("frm.xlsx"))
_m3=ct(P("frm.xlsx")); check("38 xlsx uncached formula surfaced (not silent blank)", "=A2+B2" in _m3)

# csv/tsv: native MarkItDown does NOT escape '|'/newlines or tabulate .tsv
open(P("pipe.csv"),"w",encoding="utf-8",newline="").write('Note,Detail\n"a|b","l1\nl2"\nC,"x, y"\n')
_mc=ct(P("pipe.csv")); _rows=[r for r in _mc.splitlines() if r.startswith("|")]; _cols=[r.replace("\\|","").count("|") for r in _rows]
check("39 csv escapes pipe+newline (table not broken)", "a\\|b" in _mc and "l1<br>l2" in _mc and len(set(_cols))==1, "cols=%s"%_cols)
check("40 tsv tabulated (not raw text)", "| a | b |" in ct(P("data.tsv")))
open(P("semi.csv"),"w",encoding="utf-8").write("name;city;score\nAnna;Moscow;5\n")
check("41 semicolon csv tabulated","| name | city | score |" in ct(P("semi.csv")))
open(P("crlf.csv"),"wb").write(b"h1,h2\r\nv1,v2\r\n")
_mcl=ct(P("crlf.csv")); check("42 crlf csv no stray <br>/CR", "| v1 | v2 |" in _mcl and "<br>" not in _mcl and "\r" not in _mcl)

# --- bug-hunt regression fixes (charset window, utf-16, sniff, url-name, broken pipe) ---
# 43: non-ASCII bytes start PAST the 256KB detection window -> must not mis-commit to utf-8
open(P("bigtail.csv"),"wb").write(("filler_ascii_line,value,note\n"*(262144//29+60)).encode("ascii")+"Привет,мир,конец\n".encode("cp1251"))
check("43 late single-byte tail not mis-utf8 (no false fail)", "Привет" in ct(P("bigtail.csv")))
# 44: BOM-less UTF-16 must tabulate, not become NUL-interleaved mojibake
open(P("u16le.csv"),"wb").write("Name,City\nAlice,NYC\n".encode("utf-16-le"))
_u16=ct(P("u16le.csv")); check("44 utf-16 no-BOM real table (no NUL mojibake)", "| Name | City |" in _u16 and "\x00" not in _u16)
# 45: comma CSV with a quoted ';' notes cell must stay multi-column (sniff is quoting-aware)
open(P("notes.csv"),"w",encoding="utf-8",newline="").write('"see a; b; c",Name,Score\n"x; y",Alice,10\n')
check("45 quoted-semicolon doesn't collapse comma csv", "| Name | Score |" in ct(P("notes.csv")))
# 46: URL-derived name must never contain ':' (a bare ':' becomes a hidden NTFS ADS = data loss)
check("46 url colon sanitized (no NTFS ADS)", ":" not in conv.out_name_for("http://h.example.com/report:final.pdf") and ":" not in conv.out_name_for("http://x.com:8080"))
# 47: URL fragment/query stripped, '/' in a query not mis-parsed as the path
check("47 url fragment/query stripped", conv.out_name_for("http://x.com/a.pdf#frag?y=1")=="a.pdf.md" and conv.out_name_for("http://x.com/a?b=/c/d")=="a.md")
# 48: a closed stdout pipe (BrokenPipeError) must not raise out of the echo (file is authoritative)
class _BrokenOut:
    encoding="utf-8"; buffer=None
    def write(self,*a): raise BrokenPipeError(32,"Broken pipe")
    def fileno(self): raise OSError("no fileno")
    def flush(self): pass
_saved=sys.stdout; _ok=True
try:
    sys.stdout=_BrokenOut(); conv._echo_stdout("hello world")
except BaseException:
    _ok=False
finally:
    sys.stdout=_saved
check("48 broken-pipe echo no-raise (no false exit1)", _ok)

tb=sum(b.count("Traceback (most recent call last)") for b in stderr_blobs)
check("35 ZERO tracebacks across CLI runs",tb==0,"tracebacks=%d"%tb)

passed=sum(1 for _,ok,_ in results if ok); failed=len(results)-passed
w=max(len(n) for n,_,_ in results)
print("="*70)
for n,ok,d in results: print(f"  {'PASS' if ok else 'FAIL'}  {n:<{w}}  {'' if ok else d}")
print("="*70)
print(f"TOTAL: {passed}/{len(results)} ({passed*100//len(results)}%), {failed} failed | tracebacks={tb}")
sys.exit(1 if failed else 0)
