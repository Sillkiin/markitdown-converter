#!/usr/bin/env python3
"""
selftest.py - Build small sample files for each supported format, convert them
with scripts/convert.py, and print a pass/fail summary.

Formats whose sample-builder dependency isn't installed are skipped (reported as
SKIP), not failed - so this is safe to run in any environment. A .docx that
converts via the graceful fallback still counts as a pass.

Run:
    python tests/selftest.py
"""
import os
import sys
import json
import zipfile
import tempfile
import subprocess

HERE = os.path.dirname(os.path.abspath(__file__))
CONVERT = os.path.join(HERE, "..", "scripts", "convert.py")


def build_samples(d):
    """Create one sample per format we can build here. Returns {label: path}."""
    samples = {}

    # text-y formats need no extra deps
    p = os.path.join(d, "data.json")
    json.dump({"product": "Widget", "price": 42, "tags": ["a", "b"]}, open(p, "w"))
    samples["json"] = p

    p = os.path.join(d, "table.csv")
    open(p, "w").write("name,role\nNik,founder\nIvan,dev\n")
    samples["csv"] = p

    p = os.path.join(d, "page.html")
    open(p, "w").write("<h1>Title</h1><p>Some <b>bold</b> text.</p><ul><li>a</li><li>b</li></ul>")
    samples["html"] = p

    p = os.path.join(d, "feed.xml")
    open(p, "w").write('<?xml version="1.0"?><catalog><item><name>Pen</name></item></catalog>')
    samples["xml"] = p

    p = os.path.join(d, "bundle.zip")
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("inner.csv", "a,b\n1,2\n")
    samples["zip"] = p

    # optional builders
    try:
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active
        ws.append(["Metric", "Value"]); ws.append(["MRR", "$420k"])
        p = os.path.join(d, "book.xlsx"); wb.save(p); samples["xlsx"] = p
    except Exception:
        pass

    try:
        from pptx import Presentation
        pr = Presentation(); s = pr.slides.add_slide(pr.slide_layouts[1])
        s.shapes.title.text = "Deck"; s.placeholders[1].text_frame.text = "Point one"
        p = os.path.join(d, "deck.pptx"); pr.save(p); samples["pptx"] = p
    except Exception:
        pass

    try:
        from docx import Document
        doc = Document(); doc.add_heading("Report", 1); doc.add_paragraph("Body text.")
        p = os.path.join(d, "doc.docx"); doc.save(p); samples["docx"] = p
    except Exception:
        pass

    try:
        from reportlab.pdfgen import canvas
        p = os.path.join(d, "file.pdf")
        c = canvas.Canvas(p); c.drawString(72, 720, "Hello PDF"); c.save()
        samples["pdf"] = p
    except Exception:
        pass

    return samples


def main():
    all_formats = ["pdf", "docx", "xlsx", "pptx", "csv", "html", "json", "xml", "zip"]
    with tempfile.TemporaryDirectory() as d:
        samples = build_samples(d)
        outdir = os.path.join(d, "out"); os.makedirs(outdir, exist_ok=True)
        results = []
        for fmt in all_formats:
            if fmt not in samples:
                results.append((fmt, "SKIP", "builder dependency not installed"))
                continue
            proc = subprocess.run(
                [sys.executable, CONVERT, samples[fmt], "-o", outdir + "/"],
                capture_output=True, text=True)
            note = ""
            if "NOTE:" in proc.stderr:
                note = "(fallback)"
            results.append((fmt, "PASS" if proc.returncode == 0 else "FAIL",
                            note or proc.stderr.strip().splitlines()[-1] if proc.stderr else note))

    width = max(len(f) for f, _, _ in results)
    passed = sum(1 for _, s, _ in results if s == "PASS")
    failed = sum(1 for _, s, _ in results if s == "FAIL")
    skipped = sum(1 for _, s, _ in results if s == "SKIP")
    print("MarkItDown converter self-test\n")
    for fmt, status, note in results:
        print("  %-*s  %-4s  %s" % (width, fmt, status, note))
    print("\n%d passed, %d failed, %d skipped" % (passed, failed, skipped))
    sys.exit(1 if failed else 0)


if __name__ == "__main__":
    main()
